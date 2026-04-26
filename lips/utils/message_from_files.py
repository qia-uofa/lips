import base64
import mimetypes
from pathlib import Path

# Only non-text MIME types need explicit mapping; anything unmapped is read as text.
TYPE_MAP = {
    # images
    'image/jpeg':    'image',
    'image/png':     'image',
    'image/gif':     'image',
    'image/webp':    'image',
    'image/bmp':     'image',
    'image/tiff':    'image',
    'image/svg+xml': 'image',
    'image/heic':    'image',
    'image/heif':    'image',
    'image/avif':    'image',
    # documents (text-extracted, not base64)
    'application/pdf':                                                           'document',
    'application/msword':                                                        'document',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document':   'document',
    'application/vnd.ms-excel':                                                  'document',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':         'document',
    'application/vnd.ms-powerpoint':                                             'document',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'document',
    'application/epub+zip':                                                      'document',
    # audio (OpenAI input_audio only accepts wav and mp3)
    'audio/mpeg':  'audio',
    'audio/mp3':   'audio',
    'audio/wav':   'audio',
    'audio/x-wav': 'audio',
}

AUDIO_FORMAT = {
    'audio/mpeg':  'mp3',
    'audio/mp3':   'mp3',
    'audio/wav':   'wav',
    'audio/x-wav': 'wav',
}


# --- Document text extractors -------------------------------------------------
# Each extractor takes a Path and returns a single string of plain text.
# Imports are deferred so callers only pay for libraries they actually need.

def _extract_pdf_text(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(pages)


def _extract_docx_text(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        chunks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    chunks.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                chunks.append(f"[Notes]\n{notes}")
        slides.append("\n".join(chunks))
    return "\n\n".join(slides)


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True, read_only=True)
    try:
        sheets = []
        for sheet in wb.worksheets:
            rows = [f"--- Sheet: {sheet.title} ---"]
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(cells):
                    rows.append("\t".join(cells))
            sheets.append("\n".join(rows))
        return "\n\n".join(sheets)
    finally:
        wb.close()


DOCUMENT_EXTRACTORS = {
    'application/pdf':                                                           _extract_pdf_text,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document':   _extract_docx_text,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':         _extract_xlsx_text,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': _extract_pptx_text,
}


def content_block_from_file(file, repo_path, mask=''):
    path = Path(file).resolve()
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    relative = path.relative_to(repo_path)
    display_path = Path(mask) / relative if mask else relative

    media_type, _ = mimetypes.guess_type(path)
    kind = TYPE_MAP.get(media_type, 'text')

    if kind == 'text':
        data = path.read_text(encoding='utf-8', errors='replace')
        return {
            "type": "text",
            "text": f'<file path="{display_path}">\n{data}\n</file>',
        }

    if kind == 'document':
        extractor = DOCUMENT_EXTRACTORS.get(media_type)
        if extractor is None:
            raise NotImplementedError(
                f"No text extractor for {media_type} ({path.name}). "
                "Supported: PDF, DOCX, XLSX, PPTX."
            )
        text = extractor(path)
        return {
            "type": "text",
            "text": f'<document path="{display_path}">\n{text}\n</document>',
        }

    b64 = base64.standard_b64encode(path.read_bytes()).decode('utf-8')
    url = f"data:{media_type};base64,{b64}"

    if kind == 'image':
        return [
            {"type": "text",      "text": f'<image path="{display_path}">'},
            {"type": "image_url", "image_url": {"url": url}},
            {"type": "text",      "text": "</image>"},
        ]
    if kind == 'audio':
        return {
            "type": "input_audio",
            "input_audio": {"data": b64, "format": AUDIO_FORMAT[media_type]},
        }

    raise ValueError(f"Unsupported kind: {kind!r} for file: {path.name}")


def message_from_files(files, repo_path, mask=''):
    content = []
    for file in files:
        block = content_block_from_file(file, repo_path, mask=mask)
        if isinstance(block, list):
            content.extend(block)
        else:
            content.append(block)
    return {"role": "user", "content": content}